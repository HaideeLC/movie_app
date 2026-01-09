from pptx import Presentation
from pptx.util import Inches

# 建立簡報
prs = Presentation()

# 每頁內容
slides_content = [
    ("電影訂票系統 (Movie Booking System)",
     "前端：HTML + JavaScript\n後端：Python + Flask\n資料庫：SQLite"),
    
    ("專案功能",
     "主要功能：\n- 使用者登入/登出\n- 電影列表瀏覽\n- 訂票功能（登入後）\n- 訂單號生成（UUID + 日期）\n- 訂單查詢（透過訂單編號）"),
    
    ("技術架構",
     "前端 (Frontend):\n- HTML + JavaScript\n- 表單提交與模板渲染 (Jinja2)\n\n後端 (Backend):\n- Python + Flask\n- Session 管理\n- 資料庫操作 (SQLite)"),
    
    ("資料庫 ER Diagram",
     "users               movies\n+------------+      +------------+\n| id PK      |      | id PK      |\n| username   |      | title      |\n| password   |      | showtime   |\n+------------+      +------------+\n     ^                    ^\n     |                    |\n     +--------+-----------+\n              |\n          bookings\n          +------------+\n          | id PK      |\n          | order_no   |\n          | movie_id FK|\n          | customer_name|\n          | tickets    |\n          +------------+"),
    
    ("程式模組 UML",
     "+------------------+\n|      app.py      |\n+------------------+\n| - DATABASE       |\n| - app (Flask)    |\n+------------------+\n| + init_db()      |\n| + get_db()       |\n| + close_db()     |\n| + generate_order_no() |\n| + generate_unique_order_no() |\n+------------------+\n| Routes:          |\n| - login()        |\n| - logout()       |\n| - movies()       |\n| - book(movie_id) |\n| - success(order_no) |\n| - query_order()  |\n+------------------+"),
    
    ("訂票流程圖",
     "[使用者] -> /login -> Session 設定 -> / -> 選電影 -> /book/<movie_id> -> 生成訂單號 -> 存入 bookings -> /success/<order_no> -> 查詢 /order"),
    
    ("面試加分點",
     "- 登入控制 + Session 管理\n- 訂單號唯一生成 (UUID + 日期)\n- SQL Injection 防護\n- 訂單查詢使用 JOIN 顯示電影資訊\n- 前端與後端清楚分層"),
    
    ("未來擴充",
     "- 座位剩餘數量控管\n- 訂單取消功能\n- RESTful JSON API\n- 模板繼承 + 前端美化\n- RWD 支援手機和平板")
]

# 生成每頁
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 1 = Title + Content
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# 儲存檔案
prs.save("Movie_Booking_System.pptx")
print("PPT 已生成：Movie_Booking_System.pptx")

